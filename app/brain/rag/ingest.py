import os
import uuid
from pathlib import Path

import pytesseract
from dotenv import load_dotenv

load_dotenv()

if os.getenv("TESSERACT_CMD"):
    pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD")

_embedder = None


def get_embedder():
    global _embedder
    if _embedder is None:
        from sentence_transformers import SentenceTransformer

        _embedder = SentenceTransformer("all-MiniLM-L6-v2", device="cpu")
    return _embedder


def _chunk_text(text: str, max_words: int = 375, overlap: int = 38) -> list[str]:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = words[i : i + max_words]
        if chunk:
            chunks.append(" ".join(chunk))
        i += max_words - overlap
    return chunks or [text]


def _extract_pdf(file_path: str) -> list[tuple[str, str]]:
    import fitz

    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text() or ""
        pages.append((f"pages_{i + 1}_{i + 1}", text))
    if len(pages) >= 3:
        grouped = []
        step = max(1, len(pages) // 3)
        for start in range(0, len(pages), step):
            end = min(start + step, len(pages))
            chapter_id = f"pages_{start + 1}_{end}"
            text = "\n".join(p[1] for p in pages[start:end])
            grouped.append((chapter_id, text))
        return grouped
    return pages


def _ocr_pdf(file_path: str, lang: str) -> str:
    import fitz
    import pytesseract
    from PIL import Image

    doc = fitz.open(file_path)
    parts = []
    for page in doc:
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        parts.append(pytesseract.image_to_string(img, lang=lang))
    return "\n".join(parts)


def _extract_docx(file_path: str) -> list[tuple[str, str]]:
    from docx import Document

    doc = Document(file_path)
    chapters = []
    current_id = "section_1"
    current = []
    idx = 1
    for para in doc.paragraphs:
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            if current:
                chapters.append((current_id, "\n".join(current)))
            current_id = f"section_{idx}"
            idx += 1
            current = [para.text]
        else:
            current.append(para.text)
    if current:
        chapters.append((current_id, "\n".join(current)))
    return chapters or [("section_1", "")]


def _extract_pptx(file_path: str) -> list[tuple[str, str]]:
    from pptx import Presentation

    prs = Presentation(file_path)
    chapters = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        chapters.append((f"slide_{i}", "\n".join(texts)))
    return chapters


def ingest_document(file_path: str, source_language: str = "en") -> str:
    document_id = str(uuid.uuid4())
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        chapters = _extract_pdf(file_path)
        full = "\n".join(t for _, t in chapters)
        ratio = sum(c.isalnum() for c in full) / max(len(full), 1)
        if ratio < 0.3:
            lang_map = {"en": "eng", "hi": "hin"}
            lang = lang_map.get(source_language, "eng")
            if source_language not in lang_map:
                print("[OCR] unknown source_language, defaulting to eng")
            ocr_text = _ocr_pdf(file_path, lang)
            chapters = [("pages_1_ocr", ocr_text)]
    elif ext == ".docx":
        chapters = _extract_docx(file_path)
    elif ext == ".pptx":
        chapters = _extract_pptx(file_path)
    else:
        text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        parts = [p.strip() for p in text.split("\n\n") if p.strip()]
        chapters = [(f"chapter_{i + 1}", p) for i, p in enumerate(parts)] or [
            ("chapter_1", text)
        ]

    chunks, chapter_ids = [], []
    for chapter_id, text in chapters:
        for chunk in _chunk_text(text):
            chunks.append(chunk)
            chapter_ids.append(chapter_id)

    if not chunks:
        chunks = ["(empty document)"]
        chapter_ids = ["chapter_1"]

    model = get_embedder()
    embeddings = model.encode(chunks, convert_to_numpy=True).tolist()

    import chromadb

    client = chromadb.PersistentClient(path="./chroma_data")
    collection = client.get_or_create_collection(name="documents")
    collection.add(
        ids=[f"{document_id}_{i}" for i in range(len(chunks))],
        embeddings=embeddings,
        documents=chunks,
        metadatas=[
            {"document_id": document_id, "chapter_id": chapter_ids[i]}
            for i in range(len(chunks))
        ],
    )
    return document_id
